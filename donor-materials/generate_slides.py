#!/usr/bin/env python3
"""Generate a donor-facing PowerPoint slide deck for The Center for Rural AI (CRAI)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colours ──────────────────────────────────────────────────────────────────
GREEN_DARK = RGBColor(0x1A, 0x5D, 0x1A)
GREEN_MED = RGBColor(0x4A, 0x7C, 0x2E)
GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
ACCENT = RGBColor(0xE9, 0x45, 0x60)
BLACK = RGBColor(0x1E, 0x1E, 0x1E)
GREY = RGBColor(0x64, 0x64, 0x64)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _set_slide_bg(slide, color):
    """Set the background colour of a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, color):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height):
    """Add a text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def _set_para(para, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
              font_name="Calibri", italic=False):
    """Configure a paragraph with text and formatting."""
    para.text = text
    para.alignment = alignment
    font = para.font
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    font.name = font_name
    font.italic = italic


def _add_bullet(tf, text, font_size=16, color=BLACK, level=0, space_before=Pt(4),
                space_after=Pt(2)):
    """Add a bullet point to a text frame."""
    para = tf.add_paragraph()
    para.text = text
    para.level = level
    para.space_before = space_before
    para.space_after = space_after
    font = para.font
    font.size = Pt(font_size)
    font.color.rgb = color
    font.name = "Calibri"
    return para


def _stat_box(slide, left, top, value, label):
    """Add a stat box (value + label) on the slide."""
    box_w, box_h = Inches(2.5), Inches(1.5)
    shape = _add_shape(slide, left, top, box_w, box_h, GREEN_LIGHT)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_para(tf.paragraphs[0], value, font_size=32, bold=True, color=GREEN_DARK,
              alignment=PP_ALIGN.CENTER)
    p = tf.add_paragraph()
    _set_para(p, label, font_size=11, color=GREY, alignment=PP_ALIGN.CENTER)


def build_slides(output_path="donor_slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # blank

    # ── Slide 1: Title ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, GREEN_DARK)

    tf = _add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2))
    _set_para(tf.paragraphs[0], "The Center for Rural AI", font_size=44, bold=True,
              color=WHITE, alignment=PP_ALIGN.CENTER)

    tf = _add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8))
    _set_para(tf.paragraphs[0], "Building Rural Capacity in the AI Economy",
              font_size=24, color=RGBColor(0xC8, 0xE6, 0xC8), alignment=PP_ALIGN.CENTER)

    tf = _add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6))
    _set_para(tf.paragraphs[0], "Philanthropic Partnership", font_size=18,
              color=RGBColor(0xB4, 0xD2, 0xB4), alignment=PP_ALIGN.CENTER, italic=True)

    tf = _add_textbox(slide, Inches(3), Inches(5.8), Inches(7), Inches(0.5))
    _set_para(tf.paragraphs[0], "ruralai.org", font_size=16, bold=True,
              color=WHITE, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: The Opportunity ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), GREEN_DARK)

    tf = _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
    _set_para(tf.paragraphs[0], "The Opportunity", font_size=32, bold=True,
              color=WHITE, alignment=PP_ALIGN.LEFT)

    # Stats row
    _stat_box(slide, Inches(0.8), Inches(1.5), "60M+", "Rural Americans")
    _stat_box(slide, Inches(3.8), Inches(1.5), "10%+", "Share of U.S. GDP")
    _stat_box(slide, Inches(6.8), Inches(1.5), "$5M", "Partnership Ask")
    _stat_box(slide, Inches(9.8), Inches(1.5), "24 Mo", "Timeline")

    tf = _add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.5))
    _set_para(tf.paragraphs[0],
              "The Center for Rural AI (CRAI) is a newly formed 501(c)(3) nonprofit focused "
              "on a structural economic disparity in the United States: the lagging participation "
              "of rural communities in artificial intelligence.",
              font_size=16, color=BLACK)
    _add_bullet(tf,
                "Over 60 million citizens live in rural areas contributing over 10% of U.S. GDP",
                font_size=15, color=GREY)
    _add_bullet(tf,
                "Rural ecosystems are underrepresented in high-growth tech sectors",
                font_size=15, color=GREY)
    _add_bullet(tf,
                "AI and related technologies will generate tens of trillions in value by 2030",
                font_size=15, color=GREY)
    _add_bullet(tf,
                "CRAI exists to ensure rural communities capture their share",
                font_size=15, color=GREEN_DARK)

    # ── Slide 3: About CRAI ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), GREEN_DARK)

    tf = _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
    _set_para(tf.paragraphs[0], "About CRAI", font_size=32, bold=True,
              color=WHITE, alignment=PP_ALIGN.LEFT)

    tf = _add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.2))
    _set_para(tf.paragraphs[0],
              "CRAI builds the institutional capacity rural regions need to adopt, deploy, and "
              "influence AI systems \u2014 creating economic opportunities and enhancing community outcomes.",
              font_size=16, color=BLACK)

    # Founding team
    tf = _add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.5))
    _set_para(tf.paragraphs[0], "Founding Team", font_size=22, bold=True, color=GREEN_DARK)

    people = [
        ("Andrew Aitken, Executive Director",
         "Advised the U.S. White House, Microsoft, and Capital One; served on boards of "
         "numerous technology foundations."),
        ("Adam Markham, Technology Adviser",
         "Directed AI research and systems engineering in government and critical "
         "infrastructure contexts."),
        ("Marc Nager, Adviser",
         "Co-founded Startup Weekend and led community programs at Techstars; partner at "
         "one of the only VCs focused on rural ecosystems."),
    ]
    y = Inches(3.5)
    for name, bio in people:
        tf = _add_textbox(slide, Inches(1.0), y, Inches(11), Inches(0.8))
        _set_para(tf.paragraphs[0], name, font_size=15, bold=True, color=GREEN_DARK)
        p = tf.add_paragraph()
        _set_para(p, bio, font_size=13, color=GREY)
        y += Inches(1.1)

    # ── Slide 4: Strategic Partnerships ──────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), GREEN_DARK)

    tf = _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
    _set_para(tf.paragraphs[0], "Strategic Partnerships & Policy Alignment",
              font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Left column: FLC partnership
    tf = _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5))
    _set_para(tf.paragraphs[0], "Fort Lewis College Partnership", font_size=20, bold=True,
              color=GREEN_DARK)

    tf = _add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.0))
    _set_para(tf.paragraphs[0],
              "CRAI and the AI Institute at Fort Lewis College in Durango, Colorado \u2014 "
              "a rural public institution advancing comprehensive AI education \u2014 have "
              "initiated discussions on collaborative pathways.",
              font_size=14, color=BLACK)
    _add_bullet(tf, "Direct insight into rural institutional constraints", font_size=14, color=GREY)
    _add_bullet(tf, "Engagement with frontier AI companies", font_size=14, color=GREY)
    _add_bullet(tf, "Rural use cases inform product design", font_size=14, color=GREY)

    # Right column: Federal alignment
    tf = _add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5))
    _set_para(tf.paragraphs[0], "Federal Policy Alignment", font_size=20, bold=True,
              color=GREEN_DARK)

    tf = _add_textbox(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(3.0))
    _set_para(tf.paragraphs[0],
              "Federal policy is increasingly aligned with addressing regional disparities.",
              font_size=14, color=BLACK)
    _add_bullet(tf, "CHIPS and Science Act place-based programs", font_size=14, color=GREY)
    _add_bullet(tf, "Regional Technology & Innovation Hubs (Tech Hubs)", font_size=14, color=GREY)
    _add_bullet(tf, "Up to ~$10B over five years for distributed innovation", font_size=14, color=GREY)
    _add_bullet(tf, "Structural support for rural and underserved communities", font_size=14,
                color=GREEN_DARK)

    # ── Slide 5: Programs (1/2) ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), GREEN_DARK)

    tf = _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
    _set_para(tf.paragraphs[0], "Our Programs", font_size=32, bold=True,
              color=WHITE, alignment=PP_ALIGN.LEFT)

    # Program cards - top row
    programs_top = [
        ("AI 101 \u2014 Q1/26",
         "Plain-language, hands-on half-day workshop for small business owners, "
         "municipal staff, nonprofit leaders, and farmers. No technical background required.",
         [
             "Working understanding of AI capabilities and limitations",
             "Hands-on practice with real tools",
             "Personal 30-day action plan",
         ]),
        ("AI Ignition \u2014 Q2/26",
         "Capacity-building initiative for rural higher education institutions and nonprofits.",
         [
             "AI readiness assessments",
             "High-impact, low-risk opportunity identification",
             "Focused 90-day pilots with measurable outcomes",
         ]),
    ]

    for i, (title, desc, bullets) in enumerate(programs_top):
        left = Inches(0.8 + i * 6.2)
        box = _add_shape(slide, left, Inches(1.4), Inches(5.6), Inches(2.6), GREEN_LIGHT)

        tf = _add_textbox(slide, left + Inches(0.3), Inches(1.5), Inches(5.0), Inches(0.5))
        _set_para(tf.paragraphs[0], title, font_size=18, bold=True, color=GREEN_DARK)

        tf = _add_textbox(slide, left + Inches(0.3), Inches(2.0), Inches(5.0), Inches(1.8))
        _set_para(tf.paragraphs[0], desc, font_size=12, color=BLACK)
        for b in bullets:
            _add_bullet(tf, b, font_size=12, color=GREY, space_before=Pt(2), space_after=Pt(1))

    # Program cards - bottom row
    programs_bottom = [
        ("Rural AI Innovation & Training Lab \u2014 Q3/26",
         "First-of-its-kind physical lab at Fort Lewis College Innovation Center.",
         [
             "90-day cohorts training ~15 faculty/staff",
             "Intensive in-person + remote applied pilot + capstone",
             "Each cohort influences 5,000+ students annually",
         ]),
        ("Peer Council \u2014 Q2/26",
         "Shared network and knowledge library for participating institutions.",
         [
             "Exchange tested use cases, tools, and insights",
             "Cultivate durable institutional knowledge",
             "Strengthen the rural AI ecosystem",
         ]),
    ]

    for i, (title, desc, bullets) in enumerate(programs_bottom):
        left = Inches(0.8 + i * 6.2)
        box = _add_shape(slide, left, Inches(4.4), Inches(5.6), Inches(2.6), GREEN_LIGHT)

        tf = _add_textbox(slide, left + Inches(0.3), Inches(4.5), Inches(5.0), Inches(0.5))
        _set_para(tf.paragraphs[0], title, font_size=18, bold=True, color=GREEN_DARK)

        tf = _add_textbox(slide, left + Inches(0.3), Inches(5.0), Inches(5.0), Inches(1.8))
        _set_para(tf.paragraphs[0], desc, font_size=12, color=BLACK)
        for b in bullets:
            _add_bullet(tf, b, font_size=12, color=GREY, space_before=Pt(2), space_after=Pt(1))

    # ── Slide 6: Programs (2/2) ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), GREEN_DARK)

    tf = _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
    _set_para(tf.paragraphs[0], "Our Programs (continued)", font_size=32, bold=True,
              color=WHITE, alignment=PP_ALIGN.LEFT)

    # AI in the Mountains Summit
    box = _add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(2.6), GREEN_LIGHT)
    tf = _add_textbox(slide, Inches(1.1), Inches(1.5), Inches(5.0), Inches(0.5))
    _set_para(tf.paragraphs[0], "AI in the Mountains Summit \u2014 Q4/26",
              font_size=18, bold=True, color=GREEN_DARK)
    tf = _add_textbox(slide, Inches(1.1), Inches(2.0), Inches(5.0), Inches(1.8))
    _set_para(tf.paragraphs[0],
              "Annual convening hosted in rotating rural mountain towns.",
              font_size=12, color=BLACK)
    _add_bullet(tf, "AI researchers, policymakers, corporate partners", font_size=12, color=GREY,
                space_before=Pt(2), space_after=Pt(1))
    _add_bullet(tf, "Rural stakeholders shaping AI strategy", font_size=12, color=GREY,
                space_before=Pt(2), space_after=Pt(1))
    _add_bullet(tf, "National and regional policy influence", font_size=12, color=GREY,
                space_before=Pt(2), space_after=Pt(1))

    # AI Fluency Platform
    box = _add_shape(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(2.6), GREEN_LIGHT)
    tf = _add_textbox(slide, Inches(7.3), Inches(1.5), Inches(5.0), Inches(0.5))
    _set_para(tf.paragraphs[0], "The AI Fluency Platform \u2014 Q2/26",
              font_size=18, bold=True, color=GREEN_DARK)
    tf = _add_textbox(slide, Inches(7.3), Inches(2.0), Inches(5.0), Inches(1.8))
    _set_para(tf.paragraphs[0],
              "Role-specific AI education tailored to rural institutional contexts.",
              font_size=12, color=BLACK)
    _add_bullet(tf, "Agentic content orchestration platform", font_size=12, color=GREY,
                space_before=Pt(2), space_after=Pt(1))
    _add_bullet(tf, "Contextualized learning for real operations", font_size=12, color=GREY,
                space_before=Pt(2), space_after=Pt(1))
    _add_bullet(tf, "Generates scalable IP beyond pilot phases", font_size=12, color=GREY,
                space_before=Pt(2), space_after=Pt(1))

    # AI Training Catalog callout
    box = _add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.8), Inches(2.3), GREEN_LIGHT)
    tf = _add_textbox(slide, Inches(1.1), Inches(4.6), Inches(11.2), Inches(0.5))
    _set_para(tf.paragraphs[0], "Powering the Platform: AI Training Catalog",
              font_size=18, bold=True, color=GREEN_DARK)
    tf = _add_textbox(slide, Inches(1.1), Inches(5.1), Inches(11.2), Inches(1.5))
    _set_para(tf.paragraphs[0],
              "An intelligent agent-based system that discovers free, open-source AI training "
              "content across the web and generates contextualized curricula for rural communities.",
              font_size=14, color=BLACK)
    _add_bullet(tf, "269+ curated resources across 13 AI domains (ML, NLP, Computer Vision, GenAI, and more)",
                font_size=13, color=GREY, space_before=Pt(2), space_after=Pt(1))
    _add_bullet(tf, "Rural contextualization: urban examples replaced with agriculture, healthcare, "
                "and small-town business scenarios",
                font_size=13, color=GREY, space_before=Pt(2), space_after=Pt(1))
    _add_bullet(tf, "Automated discovery, quality scoring, deduplication, and best-practices curriculum generation",
                font_size=13, color=GREY, space_before=Pt(2), space_after=Pt(1))

    # ── Slide 7: The Investment Thesis ───────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), GREEN_DARK)

    tf = _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
    _set_para(tf.paragraphs[0], "The Investment Thesis", font_size=32, bold=True,
              color=WHITE, alignment=PP_ALIGN.LEFT)

    tf = _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.5))
    _set_para(tf.paragraphs[0],
              "Rural regions possess underutilized competitive advantages:",
              font_size=18, color=BLACK)
    _add_bullet(tf, "Lower operating costs and stronger long-term community retention",
                font_size=16, color=GREY)
    _add_bullet(tf, "Rich real-world data environments in agriculture, healthcare, and logistics",
                font_size=16, color=GREY)
    _add_bullet(tf, "Highly relevant to practical AI use cases",
                font_size=16, color=GREY)

    tf = _add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(1.0))
    _set_para(tf.paragraphs[0],
              "Persistent training, infrastructure, and digital adoption gaps curtail rural "
              "communities\u2019 ability to participate at scale in the emerging AI economy.",
              font_size=16, color=GREY, italic=True)

    # Why philanthropy box
    box = _add_shape(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), GREEN_LIGHT)
    tf = _add_textbox(slide, Inches(1.1), Inches(5.1), Inches(11.0), Inches(0.5))
    _set_para(tf.paragraphs[0], "Why Philanthropic Capital Matters", font_size=18,
              bold=True, color=GREEN_DARK)
    tf = _add_textbox(slide, Inches(1.1), Inches(5.6), Inches(11.0), Inches(1.2))
    _set_para(tf.paragraphs[0],
              "Federal and corporate investment tends to follow demonstrated results; it rarely "
              "initiates them. Early philanthropic capital funds the assessments, pilots, and "
              "infrastructure that generate the evidence of impact required to unlock scalable "
              "public and private funding streams.",
              font_size=14, color=BLACK)

    # ── Slide 8: The Ask ─────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), GREEN_DARK)

    tf = _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
    _set_para(tf.paragraphs[0], "Our Invitation", font_size=32, bold=True,
              color=WHITE, alignment=PP_ALIGN.LEFT)

    # Big ask
    box = _add_shape(slide, Inches(3.5), Inches(1.5), Inches(6.3), Inches(1.2), GREEN_LIGHT)
    tf = _add_textbox(slide, Inches(3.8), Inches(1.6), Inches(5.8), Inches(1.0))
    _set_para(tf.paragraphs[0], "$5,000,000 over 24 months", font_size=28, bold=True,
              color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

    tf = _add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.5))
    _set_para(tf.paragraphs[0], "This investment will support:", font_size=18,
              color=BLACK)

    investments = [
        "AI Ignition pilots at 8\u201310 rural higher education institutions",
        "Launch of the AI Innovation and Training Lab at Fort Lewis College",
        "Development and deployment of the AI Fluency agentic curriculum platform",
        "Launch of the 2026 AI in the Mountains Summit",
        "Core operations and strategic partnerships",
    ]
    tf = _add_textbox(slide, Inches(1.0), Inches(3.8), Inches(11.0), Inches(2.5))
    _set_para(tf.paragraphs[0], investments[0], font_size=16, color=GREY)
    for inv in investments[1:]:
        _add_bullet(tf, inv, font_size=16, color=GREY)

    # ── Slide 9: Closing ─────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, GREEN_DARK)

    tf = _add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(2.0))
    _set_para(tf.paragraphs[0],
              "If this opportunity resonates, we welcome a conversation about how your "
              "partnership can help ensure rural communities lead, not follow, in the AI economy.",
              font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER, italic=True)

    tf = _add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.5))
    _set_para(tf.paragraphs[0], "Andrew Aitken", font_size=22, bold=True,
              color=WHITE, alignment=PP_ALIGN.CENTER)

    tf = _add_textbox(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.4))
    _set_para(tf.paragraphs[0], "Executive Director, The Center for Rural AI",
              font_size=16, color=RGBColor(0xC8, 0xE6, 0xC8), alignment=PP_ALIGN.CENTER)

    tf = _add_textbox(slide, Inches(1.5), Inches(5.4), Inches(10), Inches(0.4))
    _set_para(tf.paragraphs[0], "andrew@ruralai.org  |  ruralai.org",
              font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ── Save ─────────────────────────────────────────────────────────────────
    prs.save(output_path)
    print(f"Donor slides generated: {output_path}")


if __name__ == "__main__":
    build_slides()
